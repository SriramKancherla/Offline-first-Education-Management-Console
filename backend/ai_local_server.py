from fastapi import FastAPI, UploadFile, File, HTTPException
from gemma_service_local import build_prompt, generate_with_gemma
from pydantic import BaseModel
from pptx import Presentation
import os
import uuid
import json
import faiss
import numpy as np
import requests
from bs4 import BeautifulSoup
import PyPDF2
import docx
import re
import chardet
from sentence_transformers import SentenceTransformer


# =========================================================
# APP INIT
# =========================================================
app = FastAPI(title="AI Learning Assistant API")

UPLOAD_FOLDER = "uploads"
DATA_FOLDER = "data"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DATA_FOLDER, exist_ok=True)

# local embedding model (still HuggingFace but lightweight)
embed_model = SentenceTransformer("all-MiniLM-L6-v2")


# =========================================================
# REQUEST MODELS
# =========================================================
class QuestionRequest(BaseModel):
    question: str


class URLRequest(BaseModel):
    url: str


class GenerateRequest(BaseModel):
    difficulty: str


# =========================================================
# JSON STORAGE
# =========================================================
def save_json(data):
    with open(f"{DATA_FOLDER}/knowledge.json", "w") as f:
        json.dump(data, f, indent=4)


def load_json():
    path = f"{DATA_FOLDER}/knowledge.json"
    if os.path.exists(path):
        return json.load(open(path))
    return []


# =========================================================
# TEXT PROCESSING
# =========================================================
def simple_sentence_split(text):
    sentences = re.split(r'(?<=[.!?])\s+', text)
    return [s.strip() for s in sentences if s.strip()]


# =========================================================
# WEBSITE SCRAPER
# =========================================================
def scrape_website(url):

    headers = {"User-Agent": "MiniScraperBot/2.0"}

    r = requests.get(url, headers=headers, timeout=30)
    r.raise_for_status()

    soup = BeautifulSoup(r.text, "html.parser")
    texts = soup.find_all(text=True)

    visible_text = []
    for t in texts:
        if t.parent.name not in ["script", "style", "meta", "noscript"]:
            content = t.strip()
            if content:
                visible_text.append(content)

    return {
        "source_id": str(uuid.uuid4()),
        "source_type": "website",
        "source_path": url,
        "content": " ".join(visible_text)
    }


# =========================================================
# FILE LOADER
# =========================================================
def load_file(path):

    ext = os.path.splitext(path)[1].lower()
    text = ""

    try:

        # ---------- PDF ----------
        if ext == ".pdf":
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    if page.extract_text():
                        text += page.extract_text() + "\n"

        # ---------- DOCX ----------
        elif ext == ".docx":
            doc = docx.Document(path)
            parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            text = "\n".join(parts)

        # ---------- PPT ----------
        elif ext in [".pptx", ".ppt"]:
            prs = Presentation(path)
            parts = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())

                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    parts.append(cell.text.strip())

            text = "\n".join(parts)

        # ---------- TEXT ----------
        else:
            with open(path, "rb") as f:
                raw = f.read()

            encoding = chardet.detect(raw)["encoding"] or "utf-8"
            text = raw.decode(encoding, errors="ignore")

        text = re.sub(r"\s+", " ", text).strip()

        if len(text) < 20:
            text = "No readable educational content found."

    except Exception as e:
        text = f"Extraction error: {str(e)}"

    return {
        "source_id": str(uuid.uuid4()),
        "source_type": ext.replace(".", ""),
        "source_path": path,
        "content": text
    }


# =========================================================
# CHUNKING
# =========================================================
def create_chunks_from_json(records, size=800):

    chunks = []

    for rec in records:
        sentences = simple_sentence_split(rec["content"])
        chunk = ""

        for s in sentences:
            if len(chunk) + len(s) <= size:
                chunk += " " + s
            else:
                chunks.append({"text": chunk.strip()})
                chunk = s

        if chunk:
            chunks.append({"text": chunk.strip()})

    return chunks


# =========================================================
# VECTOR SEARCH
# =========================================================
def build_vector_store(chunks):

    texts = [c["text"] for c in chunks]
    embeddings = embed_model.encode(texts)

    index = faiss.IndexFlatL2(embeddings.shape[1])
    index.add(np.array(embeddings))

    return index, chunks


# =========================================================
# QA SEARCH
# =========================================================
def ask_question(question, index, chunks):

    q_emb = embed_model.encode([question])
    D, I = index.search(np.array(q_emb), k=3)

    collected = [chunks[i]["text"] for i in I[0]]
    combined = " ".join(collected)

    sentences = re.split(r'(?<=[.!?])\s+', combined)

    return "\n".join(f"- {s}" for s in sentences[:6])


# =========================================================
# ROOT
# =========================================================
@app.get("/")
def home():
    return {"message": "AI Learning Assistant API Running (LOCAL MODEL)"}


# =========================================================
# FILE UPLOAD
# =========================================================
@app.post("/upload")
async def upload_document(file: UploadFile = File(...)):

    path = os.path.join(UPLOAD_FOLDER, file.filename)

    with open(path, "wb") as f:
        f.write(await file.read())

    rec = load_file(path)

    data = load_json()
    data.append(rec)
    save_json(data)

    return {"message": "File added to knowledge base"}


# =========================================================
# SCRAPE
# =========================================================
@app.post("/scrape")
def scrape(req: URLRequest):

    try:
        rec = scrape_website(req.url)

        data = load_json()
        data.append(rec)
        save_json(data)

        return {"message": "Website added"}

    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


# =========================================================
# ASK
# =========================================================
@app.post("/ask")
def ask(req: QuestionRequest):

    data = load_json()
    if not data:
        raise HTTPException(status_code=400, detail="No knowledge available")

    chunks = create_chunks_from_json(data)
    index, stored = build_vector_store(chunks)

    answer = ask_question(req.question, index, stored)

    return {"answer": answer}


# =========================================================
# GENERATE WORKSHEET
# =========================================================
@app.post("/generate/worksheet")
def worksheet(req: GenerateRequest):

    data = load_json()
    if not data:
        raise HTTPException(status_code=400, detail="No knowledge available")

    chunks = create_chunks_from_json(data)
    combined = " ".join([c["text"] for c in chunks[:3]])

    prompt = build_prompt(combined, req.difficulty, "worksheet")
    result = generate_with_gemma(prompt)

    return {"worksheet": result}


# =========================================================
# GENERATE ASSESSMENT
# =========================================================
@app.post("/generate/assessment")
def assessment(req: GenerateRequest):

    data = load_json()
    if not data:
        raise HTTPException(status_code=400, detail="No knowledge available")

    chunks = create_chunks_from_json(data)
    combined = " ".join([c["text"] for c in chunks[:3]])

    prompt = build_prompt(combined, req.difficulty, "assessment")
    result = generate_with_gemma(prompt)

    return {"assessment": result}
