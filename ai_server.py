from fastapi import FastAPI, UploadFile, File, HTTPException
from gemma_service import build_prompt, generate_with_gemma
from pydantic import BaseModel
from pptx import Presentation
import os
import uuid
import json
import faiss
import numpy as np
import requests
from bs4 import BeautifulSoup
import PyPDF2
import docx
import re

from sentence_transformers import SentenceTransformer


app = FastAPI(title="AI Learning Assistant API")

UPLOAD_FOLDER = "uploads"
DATA_FOLDER = "data"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DATA_FOLDER, exist_ok=True)

model = SentenceTransformer("all-MiniLM-L6-v2")


class QuestionRequest(BaseModel):
    question: str


class URLRequest(BaseModel):
    url: str


class GenerateRequest(BaseModel):
    difficulty: str


def save_json(data):
    with open(f"{DATA_FOLDER}/knowledge.json", "w") as f:
        json.dump(data, f, indent=4)


def load_json():
    path = f"{DATA_FOLDER}/knowledge.json"
    if os.path.exists(path):
        return json.load(open(path))
    return []


def simple_sentence_split(text):
    sentences = re.split(r'(?<=[.!?])\s+', text)
    return [s.strip() for s in sentences if s.strip()]


def scrape_website(url):
    headers = {"User-Agent": "MiniScraperBot/2.0"}

    r = requests.get(url, headers=headers, timeout=30)
    r.raise_for_status()

    soup = BeautifulSoup(r.text, "html.parser")
    texts = soup.find_all(text=True)

    visible_text = []
    for t in texts:
        if t.parent.name not in ["script", "style", "meta", "noscript"]:
            content = t.strip()
            if content:
                visible_text.append(content)

    return {
        "source_id": str(uuid.uuid4()),
        "source_type": "website",
        "source_path": url,
        "content": " ".join(visible_text)
    }


import chardet

def load_file(path):
    ext = os.path.splitext(path)[1].lower()
    text = ""

    try:

        # ===== PDF EXTRACTION =====
        if ext == ".pdf":
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)

                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"

        # ===== DOCX EXTRACTION =====
        elif ext == ".docx":
            doc = docx.Document(path)

            parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text.strip())

            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text.strip())

            text = "\n".join(parts)

        # ===== PPTX EXTRACTION (VERY ROBUST) =====
        elif ext in [".pptx", ".ppt"]:
            prs = Presentation(path)
            parts = []

            for slide in prs.slides:
                for shape in slide.shapes:

                    # Normal text shapes
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())

                    # Tables inside slides
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    parts.append(cell.text.strip())

            text = "\n".join(parts)

        # ===== TEXT BASED FORMATS =====
        elif ext in [".txt", ".csv", ".log", ".md", ".json", ".py", ".html", ".xml", ".rtf"]:

            with open(path, "rb") as f:
                raw = f.read()

            detected = chardet.detect(raw)
            encoding = detected.get("encoding", "utf-8")

            text = raw.decode(encoding, errors="ignore")

        # ===== LAST RESORT FALLBACK =====
        else:
            with open(path, "rb") as f:
                raw = f.read()

            detected = chardet.detect(raw)
            encoding = detected.get("encoding", "utf-8")

            text = raw.decode(encoding, errors="ignore")

        # ===== CLEANUP PHASE =====
        text = text.strip()

        # Remove non-printable binary junk
        text = re.sub(r"[^\x09\x0A\x0D\x20-\x7E\x80-\xFF]", " ", text)

        # Collapse whitespace
        text = re.sub(r"\s+", " ", text)

        # FINAL VALIDATION
        if len(text) < 20:
            text = "No readable educational content could be extracted."

    except Exception as e:
        text = f"Error extracting content: {str(e)}"

    return {
        "source_id": str(uuid.uuid4()),
        "source_type": ext.replace(".", ""),
        "source_path": path,
        "content": text
    }


def create_chunks_from_json(json_records, size=800):
    chunks = []

    for rec in json_records:
        sentences = simple_sentence_split(rec["content"])
        chunk = ""

        for s in sentences:
            if len(chunk) + len(s) <= size:
                chunk += " " + s
            else:
                chunks.append({
                    "text": chunk.strip(),
                    "source_type": rec["source_type"],
                    "source_path": rec["source_path"]
                })
                chunk = s

        if chunk:
            chunks.append({
                "text": chunk.strip(),
                "source_type": rec["source_type"],
                "source_path": rec["source_path"]
            })

    return chunks


def build_vector_store(chunks):
    texts = [c["text"] for c in chunks]
    embeddings = model.encode(texts)

    index = faiss.IndexFlatL2(embeddings.shape[1])
    index.add(np.array(embeddings))

    return index, chunks


def clean_text(text):
    text = BeautifulSoup(text, "html.parser").get_text()
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def ask_question(question, index, chunks):
    q_emb = model.encode([question])
    D, I = index.search(np.array(q_emb), k=3)

    used_sources = set()
    collected_text = []

    for i in I[0]:
        chunk = chunks[i]
        cleaned = clean_text(chunk["text"])
        collected_text.append(cleaned)
        used_sources.add((chunk["source_type"], chunk["source_path"]))

    full_text = " ".join(collected_text)

    sentences = re.split(r'(?<=[.!?])\s+', full_text)

    answer = "\n".join(f"- {s}" for s in sentences[:6])

    sources = "\n".join([f"{s.upper()}: {p}" for s, p in used_sources])

    return f"""
ANSWER FROM KNOWLEDGE BASE:

{answer}

SOURCES USED:
{sources}
    """.strip()


def generate_learning_material(chunks, difficulty, mode):

    if not chunks:
        return "No knowledge available to generate material."

    combined = "\n".join([c["text"] for c in chunks[:3]])

    prompt = build_prompt(combined, difficulty, mode)

    result = generate_with_gemma(prompt)

    return result


@app.get("/")
def home():
    return {"message": "AI Learning Assistant API Running"}


@app.post("/upload")
async def upload_document(file: UploadFile = File(...)):
    path = os.path.join(UPLOAD_FOLDER, file.filename)

    with open(path, "wb") as f:
        f.write(await file.read())

    rec = load_file(path)

    data = load_json()
    data.append(rec)
    save_json(data)

    return {"message": "File added to knowledge base"}


@app.post("/scrape")
def scrape(req: URLRequest):
    try:
        rec = scrape_website(req.url)

        data = load_json()
        data.append(rec)
        save_json(data)

        return {"message": "Website added to knowledge base"}

    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


@app.post("/ask")
def ask(req: QuestionRequest):
    data = load_json()

    if not data:
        raise HTTPException(status_code=400, detail="No knowledge available")

    chunks = create_chunks_from_json(data)
    index, stored = build_vector_store(chunks)

    answer = ask_question(req.question, index, stored)

    return {"answer": answer}


@app.post("/generate/worksheet")
def worksheet(req: GenerateRequest):
    data = load_json()

    if not data:
        raise HTTPException(status_code=400, detail="No knowledge available")

    chunks = create_chunks_from_json(data)

    result = generate_learning_material(chunks, req.difficulty, "worksheet")

    return {"worksheet": result}


@app.post("/generate/assessment")
def assessment(req: GenerateRequest):
    data = load_json()

    if not data:
        raise HTTPException(status_code=400, detail="No knowledge available")

    chunks = create_chunks_from_json(data)

    result = generate_learning_material(chunks, req.difficulty, "assessment")

    return {"assessment": result}
